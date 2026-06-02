import json
import re
import logging
import os
import urllib.parse
import base64
from io import BytesIO
from typing import Any, AsyncGenerator

import httpx
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn

from pygments import highlight
from pygments.lexers import get_lexer_by_name, guess_lexer, TextLexer
from pygments.token import Token
from pygments.util import ClassNotFound

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("ppt-service")

app = FastAPI(title="PPT生成服务")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

DEEPSEEK_API_KEY = "sk-3a0cd13b507e41baa5e7d4d37fccde6c"
DEEPSEEK_BASE_URL = "https://api.deepseek.com"
DEEPSEEK_MODEL = "deepseek-v4-pro"

SPARK_API_KEY = "6785c837a76d654f2dd5148957c53c78"
SPARK_BASE_URL = "https://spark-api-open.xf-yun.com/v1"
SPARK_MODEL = "4.0Ultra"

DOUBAO_API_KEY = "ark-03a24903-14ef-49a0-9c6c-89e75bd8c7d8-df871"
DOUBAO_BASE_URL = "https://ark.cn-beijing.volces.com/api/v3"
DOUBAO_IMAGE_MODEL = "doubao-seedream-4-0-250828"

XF_SEARCH_API_PASSWORD = "iskDYCaZCkeJUhBfoQtf:PpOUWSuaXMJmyKxHihpZ"  # 从 https://console.xfyun.cn/services/cbm 获取，不是 Spark 的 api-key
XF_SEARCH_URL = "https://search-api-open.cn-huabei-1.xf-yun.com/v2/search"

DOWNLOAD_DATA_DIR = "D:/OTHERS/learnAi/backend/downloadData"

DOC_TYPE_DIR = {
    "ppt": "ppt",
    "lecture": "knowledge",
    "exercise": "exercise",
    "reading": "exploreReading",
}

SLIDE_GENERATION_PROMPT = """你是一个专业的PPT内容生成专家。请根据用户输入的课程信息，生成一份完整的教学PPT内容。

严格按照以下JSON格式输出，不要输出任何其他内容：

{
  "slides": [
    {
      "type": "title",
      "title": "PPT主标题",
      "subtitle": "副标题",
      "author": "作者"
    },
    {
      "type": "content",
      "title": "页面标题",
      "bullets": ["要点1", "要点2", "要点3"]
    },
    {
      "type": "content",
      "title": "页面标题",
      "bullets": ["要点1", "要点2"]
    },
    {
      "type": "table",
      "title": "表格页标题",
      "headers": ["列头1", "列头2"],
      "rows": [
        ["数据1", "数据2"],
        ["数据3", "数据4"]
      ]
    },
    {
      "type": "end",
      "title": "谢谢观看",
      "subtitle": "Q&A"
    }
  ]
}

要求：
1. 包含1个title类型封面页 + 4-8个content类型内容页 + 0-1个table类型表格页 + 1个end类型结束页
2. content页面每页3-5个要点，每个要点不超过30字
3. table页面2-4行数据
4. 所有标题简洁有力，不超过20字
5. 内容专业准确，适合教学场景
6. 只输出JSON，不要有任何解释性文字"""

LECTURE_PROMPT = """你是一个专业的教育内容生成专家。请根据用户输入的专业背景、课程内容、知识短板和学习需求，生成一份结构化的知识讲解文档。

文档使用Markdown格式输出，包含以下结构：

# 文档标题

## 一、知识概述
简要介绍知识点的背景和核心概念，2-3段。

## 二、核心概念详解
分点详细讲解核心概念，每点独立一段，可使用有序列表。

## 三、代码示例（如适用）
如果知识点涉及编程，请提供带注释的代码示例，使用```语言名 代码块```格式。
如果知识点不涉及编程，此部分改为"案例分析"。

## 四、常见误区与注意事项
列出3-4个常见误区，使用无序列表。

## 五、练习建议
给出针对性的学习建议和练习方向。

## 六、总结
用简洁的语言总结本章重点，2-3句话。

要求：
1. 内容专业准确，深入浅出
2. 总字数800-2000字
3. Markdown格式规范，代码块标注语言类型
4. 重点针对用户的知识短板展开"""

EXERCISE_PROMPT = """你是一个专业的出题专家。请根据用户输入的专业背景、课程内容、知识短板和学习需求，生成一套练习题。

文档使用Markdown格式输出，包含以下结构：

# 练习题

## 一、选择题（5题）
每题4个选项，标注正确答案。格式：
1. 题目
   A. 选项A
   B. 选项B
   C. 选项C
   D. 选项D
   答案：X

## 二、填空题（3题）
每题留空用____表示，标注答案。

## 三、简答题（2题）
考察概念理解和应用能力。

## 四、编程题（1题，如适用）
如有编程背景，提供一道编程练习，附参考代码。
代码格式：```语言名 代码块```

## 五、答案与解析

要求：
1. 题目难度适中，覆盖知识短板
2. 选择题干扰项要有迷惑性但逻辑合理
3. 答案和解析要详细"""

EXERCISE_SUMMARY_PROMPT = """你是一个专业的出题专家。用户通过搜索引擎获取了一批与该知识点相关的网页资料。

请你结合搜索结果和用户需求，生成一套练习题：

# 练习题

## 一、选择题（5题）
每题4个选项，标注正确答案。格式：
1. 题目
   A. 选项A
   B. 选项B
   C. 选项C
   D. 选项D
   答案：X

## 二、填空题（3题）
每题留空用____表示，标注答案。

## 三、简答题（2题）
考察概念理解和应用能力。

## 四、编程题（1题，如适用）
如有编程背景，提供一道编程练习，附参考代码。
代码格式：```语言名 代码块```

## 五、答案与解析

要求：
1. 参考搜索到的网页资料出题，确保题目紧扣知识点
2. 题目难度适中，覆盖用户的知识短板
3. 选择题干扰项要有迷惑性但逻辑合理
4. 答案和解析要详细"""

READING_PROMPT = """你是一个专业的阅读推荐专家。请根据用户输入的专业背景、课程内容、知识短板和学习需求，生成一份拓展阅读材料推荐。

文档使用Markdown格式输出，包含以下结构：

# 拓展阅读材料

## 一、推荐书籍（3本）
每本包含书名、作者、推荐理由、适合人群。

## 二、推荐在线资源（3-4个）
包含资源名称、网址/平台、内容简介。

## 三、推荐论文或技术文章（2篇）
适合进阶学习，简要介绍核心观点。

## 四、学习路线图
根据用户的知识短板和学习需求，给出一个分阶段的学习路线。

## 五、阅读建议
给出具体的学习方法和时间安排建议。

要求：
1. 推荐材料专业且实用
2. 针对用户的知识短板给出针对性建议
3. 学习路线切实可行"""

READING_SUMMARY_PROMPT = """你是一个专业的学习资料整理专家。用户通过搜索引擎获取了一批与该知识点相关的网页链接和摘要。

请你根据以下格式，整理一份「拓展阅读资料」文档：

# 拓展阅读资料

## 一、搜索概览
简要介绍本次搜索的主题和找到的资料概况（2-3句话）。

## 二、精选资料（选取最相关的5-8条）
每条格式：
### 资料标题
- **来源**：来源网站名称
- **摘要**：用你自己的话总结该资料的核心内容（50-80字）
- **链接**：[原文链接](URL)
- **推荐理由**：为什么推荐阅读这篇（1句话）

如果搜索结果中包含代码、算法等技术内容，请重点推荐。

## 三、核心知识点整理
根据搜索结果，提炼出3-5个核心知识点，每个用一段话简要说明。

## 四、学习建议
结合用户的知识短板，给出针对性学习建议和阅读顺序。

要求：
1. 精选的资料必须来源真实、内容可靠
2. 如果搜索结果少于3条，就只列出已有的并诚实说明
3. 核心知识点要准确、专业
4. Markdown格式规范"""


DOC_TYPE_PROMPTS = {
    "lecture": LECTURE_PROMPT,
    "exercise": EXERCISE_PROMPT,
    "reading": READING_PROMPT,
}

DOC_TYPE_LABELS = {
    "lecture": "知识讲解",
    "exercise": "练习题目",
    "reading": "拓展阅读",
    "ppt": "PPT课件",
}


class PptRequest(BaseModel):
    major: str = ""
    courseName: str = ""
    knowledgeGaps: str = ""
    learningNeeds: str = ""
    model: str = "deepseek"
    docType: str = "lecture"
    imageModel: str = "seedream"
    searchModel: str = "xfsearch"
    searchResults: str = ""
    searchItemCount: str = ""


@app.post("/api/ppt/generate")
async def generate_ppt(req: PptRequest):
    prompt_text = (f"专业：{req.major or ''}，课程内容：{req.courseName or ''}，"
                   f"知识短板：{req.knowledgeGaps or ''}，学习需求：{req.learningNeeds or ''}")
    logger.info("PPT生成请求: model=%s, prompt=%s", req.model, prompt_text)

    if req.model == "spark":
        slides_data = await call_spark(prompt_text)
    else:
        slides_data = await call_deepseek(prompt_text)

    pptx_bytes = build_pptx(slides_data, req.courseName, req.learningNeeds)

    model_tag = "spark" if req.model == "spark" else "deepseek"
    safe_name = re.sub(r'[^\w\u4e00-\u9fff]', '_', req.courseName or '课件')
    filename = f"{safe_name}_课件_{model_tag}.pptx"
    encoded = urllib.parse.quote(filename)

    return StreamingResponse(
        BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{encoded}",
        },
    )


@app.post("/api/generate/stream")
async def generate_stream(req: PptRequest):
    doc_type = req.docType or "lecture"

    prompt_text = (f"专业：{req.major or ''}，课程内容：{req.courseName or ''}，"
                   f"知识短板：{req.knowledgeGaps or ''}，学习需求：{req.learningNeeds or ''}")

    model_label = "讯飞星火" if req.model == "spark" else "DeepSeek"
    logger.info("[STREAM] 请求: docType=%s, model=%s", doc_type, model_label)

    async def event_stream() -> AsyncGenerator[dict, None]:
        file_bytes: bytes | None = None
        filename: str = ""
        mime_type: str = ""

        try:
            if doc_type == "reading":
                search_label = "讯飞ONE SEARCH" if req.searchModel == "xfsearch" else "DeepSeek"
                yield {"step": "searching", "message": f"正在调用 {search_label} 全网搜索资料...", "icon": "ri-search-line"}
                logger.info("[STREAM] ▶ 搜索: model=%s", req.searchModel)

                if req.searchModel == "xfsearch":
                    search_results = await xf_search(req.knowledgeGaps or req.courseName or "学习资料")
                else:
                    search_results = await deepseek_search(req.knowledgeGaps or req.courseName or "学习资料")

                count = len(search_results)
                if count > 0:
                    yield {"step": "search_done", "message": f"搜索到 {count} 条资料，正在调用 {model_label} 整理...", "icon": "ri-file-list-3-line"}
                    logger.info("[STREAM] ✓ 搜索完成: %d 条", count)
                else:
                    yield {"step": "search_done", "message": f"未搜索到资料，由 {model_label} 基于知识生成...", "icon": "ri-file-list-3-line"}
                    logger.info("[STREAM] - 搜索无结果")
                yield {"step": "ai_writing", "message": f"正在调用 {model_label} 整理搜索结果...", "icon": "ri-brain-line"}
                logger.info("[STREAM] ▶ 开始 AI 整理")
                ai_text = await generate_reading_doc(prompt_text, req, search_results)
            elif doc_type == "exercise":
                search_label = "讯飞ONE SEARCH" if req.searchModel == "xfsearch" else "DeepSeek"
                yield {"step": "searching", "message": f"正在调用 {search_label} 全网搜索资料...", "icon": "ri-search-line"}
                logger.info("[STREAM] ▶ 搜索: model=%s", req.searchModel)

                search_results = await xf_search(req.knowledgeGaps or req.courseName or "学习资料")

                count = len(search_results)
                if count > 0:
                    yield {"step": "search_done", "message": f"搜索到 {count} 条资料，正在调用 {model_label} 出题...", "icon": "ri-file-list-3-line"}
                    logger.info("[STREAM] ✓ 搜索完成: %d 条", count)
                else:
                    yield {"step": "search_done", "message": f"未搜索到资料，由 {model_label} 基于知识出题...", "icon": "ri-file-list-3-line"}
                    logger.info("[STREAM] - 搜索无结果")
                yield {"step": "ai_writing", "message": f"正在调用 {model_label} 基于资料出题...", "icon": "ri-brain-line"}
                logger.info("[STREAM] ▶ 开始 AI 出题")
                ai_text = await generate_exercise_doc(prompt_text, req, search_results)
            else:
                label = DOC_TYPE_LABELS.get(doc_type, doc_type)
                yield {"step": "ai_writing", "message": f"正在调用 {model_label} 生成{label}文本...", "icon": "ri-brain-line"}
                logger.info("[STREAM] ▶ 开始 AI 文本: type=%s, model=%s", doc_type, model_label)
                if doc_type == "ppt":
                    ppt_prompt = SLIDE_GENERATION_PROMPT
                    current_prompt = prompt_text
                    if hasattr(req, 'searchResults') and req.searchResults:
                        ppt_prompt = SLIDE_GENERATION_PROMPT.replace(
                            '"type": "end"',
                            '"type": "search_results","title":"📚 搜索资源","items":["资源标题 / 摘要"]},\n    {"type": "end"'
                        )
                        current_prompt += "\n\n以下是通过搜索引擎找到的相关资料，请参考并在PPT中体现：\n" + req.searchResults[:3000]
                    ai_text = await call_ai_text(ppt_prompt, current_prompt, req.model)
                else:
                    ai_text = await call_ai_text(DOC_TYPE_PROMPTS[doc_type], prompt_text, req.model)

            if not ai_text or len(ai_text) < 50:
                yield {"step": "error", "message": f"{model_label} 返回内容过短，请重试", "icon": "ri-error-warning-line"}
                return

            yield {"step": "ai_done", "message": f"{model_label} 内容生成完成 ({len(ai_text)}字)，开始生成配图...", "icon": "ri-image-add-line"}
            logger.info("[STREAM] ✓ AI文本: %d 字", len(ai_text))

            if doc_type == "ppt":
                topic = req.courseName or req.knowledgeGaps or "知识讲解"
                img_label = "豆包 Seedream 4.0" if req.imageModel == "seedream" else req.imageModel
                yield {"step": "image_gen", "message": f"正在调用 {img_label} 生成 PPT 配图...", "icon": "ri-image-edit-line"}
                logger.info("[STREAM] ▶ 开始豆包PPT配图: topic=%s", topic)
                ppt_images = await generate_ppt_images(topic)
                if ppt_images:
                    yield {"step": "image_done", "message": f"豆包PPT配图生成完成 ({len(ppt_images)}张)", "icon": "ri-check-line"}
                    logger.info("[STREAM] ✓ 豆包PPT配图: %d张", len(ppt_images))
                else:
                    yield {"step": "image_done", "message": "豆包 API Key 未配置，跳过PPT配图", "icon": "ri-information-line"}
                    logger.info("[STREAM] - 豆包PPT配图跳过")

                yield {"step": "rendering", "message": "正在使用 python-pptx 渲染 PPT 课件...", "icon": "ri-slideshow-3-line"}
                logger.info("[STREAM] ▶ 开始 python-pptx 渲染")
                obj = json.loads(extract_json(ai_text)) if isinstance(ai_text, str) else ai_text
                slides_data = obj.get("slides", [])
                # 注入搜索资源页
                if hasattr(req, 'searchResults') and req.searchResults:
                    slides_data = _inject_search_slide(slides_data, req.searchResults, req.searchItemCount)
                file_bytes = build_pptx(slides_data, req.courseName, req.learningNeeds, ppt_images)
                safe_name = re.sub(r'[^\w\u4e00-\u9fff]', '_', req.courseName or '课件')
                model_tag = "spark" if req.model == "spark" else "deepseek"
                filename = f"{safe_name}_课件_{model_tag}.pptx"
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                logger.info("[STREAM] ✓ PPT: %d bytes", len(file_bytes))
            else:
                topic = req.courseName or req.knowledgeGaps or "知识讲解"
                img_label = "豆包 Seedream 4.0" if req.imageModel == "seedream" else req.imageModel
                yield {"step": "image_gen", "message": f"正在调用 {img_label} 生成教育配图...", "icon": "ri-image-edit-line"}
                logger.info("[STREAM] ▶ 开始豆包图片: topic=%s", topic)
                image_bytes = await generate_doc_image(topic)
                if image_bytes:
                    yield {"step": "image_done", "message": f"豆包配图生成完成 ({len(image_bytes)//1024}KB)", "icon": "ri-check-line"}
                    logger.info("[STREAM] ✓ 豆包图片: %d bytes", len(image_bytes))
                else:
                    yield {"step": "image_done", "message": "豆包 API Key 未配置，跳过配图生成", "icon": "ri-information-line"}
                    logger.info("[STREAM] - 豆包跳过")
                label = DOC_TYPE_LABELS.get(doc_type, doc_type)
                yield {"step": "rendering", "message": f"正在使用 python-docx + pygments 渲染{label}文档...", "icon": "ri-file-word-2-line"}
                logger.info("[STREAM] ▶ 开始 python-docx 渲染")
                file_bytes = build_docx(ai_text, req.courseName or "文档标题", label, image_bytes)
                safe_name = re.sub(r'[^\w\u4e00-\u9fff]', '_', req.courseName or '文档')
                model_tag = "spark" if req.model == "spark" else "deepseek"
                filename = f"{safe_name}_{label}_{model_tag}.docx"
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                logger.info("[STREAM] ✓ DOCX: %d bytes", len(file_bytes))

            download_url = _save_to_disk(file_bytes, filename, doc_type)
            logger.info("[STREAM] ✓ 文件已保存: %s", download_url)

            yield {"step": "rendering_done", "message": "文件渲染完成，准备下载...", "icon": "ri-download-cloud-line"}
            yield {
                "step": "done",
                "message": "生成完成",
                "filename": filename,
                "mimeType": mime_type,
                "downloadUrl": download_url,
                "base64": base64.b64encode(file_bytes).decode("ascii"),
            }
            logger.info("[STREAM] ✓ 全部完成: %s (%d bytes)", filename, len(file_bytes))

        except Exception as e:
            logger.exception("[STREAM] 生成失败: %s", str(e))
            yield {"step": "error", "message": f"生成失败: {str(e)}", "icon": "ri-error-warning-line"}

    async def sse_gen():
        seen = set()
        async for event in event_stream():
            if event.get("step") in seen and event.get("step") != "done":
                continue
            seen.add(event.get("step"))
            yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(
        sse_gen(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/api/doc/generate")
async def generate_doc(req: PptRequest):
    doc_type = req.docType or "lecture"
    if doc_type not in DOC_TYPE_PROMPTS:
        raise HTTPException(status_code=400, detail=f"无效的文档类型: {doc_type}")

    prompt_text = (f"专业：{req.major or ''}，课程内容：{req.courseName or ''}，"
                   f"知识短板：{req.knowledgeGaps or ''}，学习需求：{req.learningNeeds or ''}")
    logger.info("文档生成请求: type=%s, model=%s", doc_type, prompt_text)

    if doc_type == "reading":
        ai_text = await generate_reading_doc(prompt_text, req)
    else:
        ai_text = await call_ai_text(DOC_TYPE_PROMPTS[doc_type], prompt_text, req.model)

    topic = req.courseName or req.knowledgeGaps or "知识讲解"
    image_bytes = await generate_doc_image(topic)

    safe_name = re.sub(r'[^\w\u4e00-\u9fff]', '_', req.courseName or '文档')
    label = DOC_TYPE_LABELS.get(doc_type, doc_type)
    model_tag = "spark" if req.model == "spark" else "deepseek"
    filename = f"{safe_name}_{label}_{model_tag}.docx"
    encoded = urllib.parse.quote(filename)

    docx_bytes = build_docx(ai_text, req.courseName or "文档标题", label, image_bytes)

    return StreamingResponse(
        BytesIO(docx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{encoded}",
        },
    )


async def generate_reading_doc(prompt_text: str, req: PptRequest,
                               search_results: list[dict[str, str]] | None = None) -> str:
    if search_results is None:
        query = req.knowledgeGaps or req.courseName or req.learningNeeds or "学习资料"
        search_results = await xf_search(query)

    if search_results:
        results_text = "以下是从全网搜索到的相关资料：\n\n"
        for i, item in enumerate(search_results):
            results_text += (
                f"{i + 1}. **{item['name']}**\n"
                f"   来源：{item.get('source', '未知')}\n"
                f"   链接：{item['url']}\n"
                f"   摘要：{item['summary'][:200]}\n\n"
            )
        user_content = (
            f"用户需求：{prompt_text}\n\n"
            f"{results_text}\n"
            f"请根据以上搜索结果为用户整理一份结构化的拓展阅读资料文档。"
        )
    else:
        user_content = (
            f"用户需求：{prompt_text}\n\n"
            f"（注意：全网搜索未返回结果，请基于你的知识为用户推荐学习资料，"
            f"但不要在文档中虚构搜索链接。）"
        )

    return await call_ai_text(READING_SUMMARY_PROMPT, user_content, req.model)


async def generate_exercise_doc(prompt_text: str, req: PptRequest,
                                search_results: list[dict[str, str]]) -> str:
    if search_results:
        results_text = "以下是从全网搜索到的相关资料：\n\n"
        for i, item in enumerate(search_results):
            results_text += (
                f"{i + 1}. **{item['name']}**\n"
                f"   来源：{item.get('source', '未知')}\n"
                f"   链接：{item['url']}\n"
                f"   摘要：{item['summary'][:200]}\n\n"
            )
        user_content = (
            f"用户需求：{prompt_text}\n\n"
            f"{results_text}\n"
            f"请根据以上搜索结果，参考其中的知识点，为用户出一套练习题。"
        )
    else:
        user_content = (
            f"用户需求：{prompt_text}\n\n"
            f"（注意：全网搜索未返回结果，请基于你的知识为用户出题。）"
        )

    return await call_ai_text(EXERCISE_SUMMARY_PROMPT, user_content, req.model)


async def xf_search(query: str, limit: int = 8) -> list[dict[str, str]]:
    if not XF_SEARCH_API_PASSWORD or XF_SEARCH_API_PASSWORD.startswith("你的"):
        logger.warning("[SEARCH] APIPassword 未配置，请在控制台 https://console.xfyun.cn/services/cbm 获取")
        return []

    body = {
        "search_params": {
            "query": query,
            "limit": limit,
            "enhance": {
                "open_full_text": True,
                "open_rerank": True,
            },
        }
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {XF_SEARCH_API_PASSWORD}",
    }

    logger.info("[SEARCH] URL=%s", XF_SEARCH_URL)
    logger.info("[SEARCH] query=%s, limit=%d", query, limit)
    logger.info("[SEARCH] APIPassword 前缀=%s...", XF_SEARCH_API_PASSWORD[:6])

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(XF_SEARCH_URL, json=body, headers=headers)

        logger.info("[SEARCH] HTTP status=%d", resp.status_code)
        if resp.status_code != 200:
            logger.error("[SEARCH] HTTP 错误: %s, body=%s", resp.status_code, resp.text[:500])
            return []

        data = resp.json()
        logger.info("[SEARCH] success=%s, err_code=%s, sid=%s",
                    data.get("success"), data.get("err_code"), data.get("sid"))

        if not data.get("success"):
            logger.error("[SEARCH] 业务失败: err_code=%s, message=%s",
                         data.get("err_code"), data.get("message", "无详细信息"))
            if data.get("err_code") == "11200":
                logger.error(
                    "[SEARCH] 授权错误 - APIPassword 可能无效，请前往 "
                    "https://console.xfyun.cn/services/cbm 确认已开通 ONE SEARCH 服务 "
                    "并复制正确的 APIPassword（不是 Spark api-key）"
                )
            return []

        documents = data.get("data", {}).get("search_results", {}).get("documents", [])
        results = []
        for doc in documents:
            results.append({
                "name": doc.get("name", "无标题"),
                "url": doc.get("url", ""),
                "summary": doc.get("summary", "") or doc.get("content", "") or "无摘要",
                "source": _extract_domain(doc.get("url", "")),
            })

        logger.info("[SEARCH] ✓ 成功: %d 条结果, query=%s", len(results), query)
        return results

    except Exception as e:
        logger.exception("[SEARCH] 异常: %s", str(e))
        return []


async def deepseek_search(query: str, limit: int = 6) -> list[dict[str, str]]:
    try:
        prompt = f"请对以下主题进行知识搜索，列出{limit}条相关概念或要点，每条用一句话概括：{query}"
        content = await call_ai_text("你是专业的知识检索助手，回复简洁准确，只列要点。", prompt, "deepseek")
        results = []
        for line in content.split('\n'):
            line = line.strip()
            if not line:
                continue
            line = re.sub(r'^\d+[\.\)、]\s*', '', line)
            line = re.sub(r'^[-*•]\s*', '', line)
            if len(line) > 10:
                results.append({
                    "name": line[:80],
                    "url": "",
                    "summary": line[:200],
                    "source": "DeepSeek AI",
                })
        logger.info("[DEEPSEEK-SEARCH] ✓ 返回 %d 条要点, query=%s", len(results), query)
        return results[:limit]
    except Exception as e:
        logger.error("[DEEPSEEK-SEARCH] 异常: %s", str(e))
        return []


def _extract_domain(url: str) -> str:
    import re as _re
    m = _re.search(r'https?://([^/]+)', url)
    return m.group(1) if m else "未知来源"


def _save_to_disk(file_bytes: bytes, filename: str, doc_type: str) -> str:
    subdir = DOC_TYPE_DIR.get(doc_type, "other")
    dir_path = os.path.join(DOWNLOAD_DATA_DIR, subdir)
    os.makedirs(dir_path, exist_ok=True)
    file_path = os.path.join(dir_path, filename)
    with open(file_path, "wb") as f:
        f.write(file_bytes)
    logger.info("[DISK] 写入: %s (%d bytes)", file_path, len(file_bytes))
    return f"/api/resource/download/{subdir}/{filename}"


async def call_ai_text(system_prompt: str, user_prompt: str, model: str) -> str:
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]
    if model == "spark":
        api_key = SPARK_API_KEY
        api_url = f"{SPARK_BASE_URL}/chat/completions"
        model_name = SPARK_MODEL
    else:
        api_key = DEEPSEEK_API_KEY
        api_url = f"{DEEPSEEK_BASE_URL}/chat/completions"
        model_name = DEEPSEEK_MODEL

    body = {
        "model": model_name,
        "messages": messages,
        "stream": False,
        "temperature": 0.7,
        "max_tokens": 8192,
    }
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}",
    }
    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(api_url, json=body, headers=headers)

    if resp.status_code != 200:
        logger.error("AI文本生成错误: %s %s", resp.status_code, resp.text[:500])
        raise HTTPException(status_code=502, detail=f"AI服务请求失败: {resp.status_code}")

    data = resp.json()
    content = data["choices"][0]["message"]["content"]
    logger.info("AI文本响应(前300字): %s", content[:300])
    return content


async def generate_doc_image(topic: str) -> bytes | None:
    if not DOUBAO_API_KEY or DOUBAO_API_KEY.startswith("这里填写"):
        logger.warning("豆包 API Key 未配置，跳过图片生成")
        return None

    return await _call_doubao_image(f"教育插图，{topic}，知识讲解配图，"
        "简洁专业的插画风格，柔和色调，浅色背景，"
        "适合教学文档使用，信息图表风格，清晰易懂")


async def generate_ppt_images(topic: str) -> list[bytes]:
    images: list[bytes] = []
    if not DOUBAO_API_KEY or DOUBAO_API_KEY.startswith("这里填写"):
        logger.warning("豆包 API Key 未配置，跳过PPT配图生成")
        return images

    subtopics = [
        f"教育幻灯片配图，{topic}，概念总览图，教学图表风格，简洁专业，柔和色调",
        f"教育幻灯片配图，{topic}，核心知识拆解图，信息图表风格，清晰易懂，浅色背景",
        f"教育幻灯片配图，{topic}，应用场景示例图，插画风格，适合教学展示",
    ]

    for idx, prompt_text in enumerate(subtopics):
        img_bytes = await _call_doubao_image(prompt_text)
        if img_bytes:
            images.append(img_bytes)
            logger.info("[PPT-IMG] ✓ 第%d张配图: %d bytes", idx + 1, len(img_bytes))
        else:
            logger.info("[PPT-IMG] - 第%d张配图生成失败，跳过", idx + 1)

    logger.info("[PPT-IMG] 共生成 %d 张配图", len(images))
    return images


async def _call_doubao_image(prompt_text: str) -> bytes | None:

    body = {
        "model": DOUBAO_IMAGE_MODEL,
        "prompt": prompt_text,
        "size": "2K",
        "response_format": "url",
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {DOUBAO_API_KEY}",
    }

    try:
        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(
                f"{DOUBAO_BASE_URL}/images/generations",
                json=body,
                headers=headers,
            )
        if resp.status_code != 200:
            logger.error("豆包图片生成失败: %s %s", resp.status_code, resp.text[:500])
            return None

        data = resp.json()
        image_url = data["data"][0]["url"]
        logger.info("豆包图片URL: %s", image_url)

        async with httpx.AsyncClient(timeout=60) as client:
            img_resp = await client.get(image_url)
        if img_resp.status_code != 200:
            logger.error("下载豆包图片失败: %s", img_resp.status_code)
            return None

        logger.info("豆包图片下载成功, size=%d bytes", len(img_resp.content))
        return img_resp.content
    except Exception as e:
        logger.error("豆包图片生成异常: %s", str(e))
        return None


def build_docx(markdown_text: str, title: str, doc_label: str, image_bytes: bytes | None = None) -> bytes:
    doc = Document()

    style = doc.styles['Normal']
    style.font.name = 'Microsoft YaHei'
    style.font.size = DocPt(11)
    style.paragraph_format.line_spacing = 1.5
    style.paragraph_format.space_after = DocPt(6)

    for level in range(1, 4):
        heading_style = doc.styles[f'Heading {level}']
        heading_style.font.name = 'Microsoft YaHei'
        heading_style.font.color.rgb = DocRGBColor(0x1A, 0x1A, 0x2E)

    code_font = 'Consolas'
    code_size = DocPt(9)
    code_bg = DocRGBColor(0xF5, 0xF5, 0xF5)

    in_code_block = False
    code_lines = []
    code_lang = 'text'
    image_inserted = False

    for raw_line in markdown_text.split('\n'):
        line = raw_line.rstrip()

        if line.startswith('```'):
            if in_code_block:
                _render_code_block(doc, code_lines, code_lang, code_font, code_size, code_bg)
                code_lines = []
                code_lang = 'text'
                in_code_block = False
            else:
                in_code_block = True
                lang = line[3:].strip()
                if lang:
                    code_lang = lang
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        if line.startswith('# '):
            doc.add_heading(line[2:].strip(), level=1)
            if image_bytes and not image_inserted:
                image_inserted = True
                try:
                    img_stream = BytesIO(image_bytes)
                    pic = doc.add_picture(img_stream, width=DocInches(5.5))
                    pic.alignment = WD_ALIGN_PARAGRAPH.CENTER
                except Exception as e:
                    logger.warning("插入配图失败: %s", str(e))
        elif line.startswith('## '):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith('- ') or line.startswith('* '):
            p = doc.add_paragraph(line[2:].strip(), style='List Bullet')
        elif line.startswith('+ '):
            p = doc.add_paragraph(line[2:].strip(), style='List Bullet')
        elif re.match(r'^\d+[\.\)]\s', line):
            p = doc.add_paragraph(re.sub(r'^\d+[\.\)]\s*', '', line), style='List Number')
        elif line.strip() == '':
            continue
        elif line.startswith('---'):
            doc.add_paragraph('─' * 40)
        else:
            p = doc.add_paragraph(line)

    if in_code_block and code_lines:
        _render_code_block(doc, code_lines, code_lang, code_font, code_size, code_bg)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _render_code_block(doc, lines, lang, font, size, bg_color):
    if not lines:
        return
    code_text = '\n'.join(lines)

    try:
        lexer = get_lexer_by_name(lang, stripall=True)
    except ClassNotFound:
        lexer = TextLexer()

    tokens = list(lexer.get_tokens(code_text))

    para = doc.add_paragraph()
    para.paragraph_format.space_before = DocPt(6)
    para.paragraph_format.space_after = DocPt(6)

    pPr = para._element.get_or_add_pPr()
    shd = pPr.makeelement(qn('w:shd'), {
        qn('w:val'): 'clear',
        qn('w:color'): 'auto',
        qn('w:fill'): 'F5F5F5',
    })
    pPr.insert(0, shd)

    for token_type, text in tokens:
        run = para.add_run(text)
        run.font.name = font
        run.font.size = size
        run.font.color.rgb = _pygments_color(token_type)

    para.add_run('')


def _pygments_color(token_type) -> DocRGBColor:
    color_map = {
        Token.Keyword: DocRGBColor(0x00, 0x00, 0xFF),
        Token.Keyword.Constant: DocRGBColor(0x00, 0x00, 0xFF),
        Token.Keyword.Declaration: DocRGBColor(0x00, 0x00, 0xFF),
        Token.Keyword.Namespace: DocRGBColor(0x00, 0x00, 0xFF),
        Token.Keyword.Type: DocRGBColor(0x2B, 0x91, 0xAF),
        Token.Name.Class: DocRGBColor(0x2B, 0x91, 0xAF),
        Token.Name.Function: DocRGBColor(0x79, 0x56, 0x1D),
        Token.Name.Builtin: DocRGBColor(0x00, 0x80, 0x00),
        Token.String: DocRGBColor(0xA3, 0x15, 0x15),
        Token.String.Doc: DocRGBColor(0x00, 0x80, 0x00),
        Token.Number: DocRGBColor(0x09, 0x85, 0x58),
        Token.Comment: DocRGBColor(0x00, 0x80, 0x00),
        Token.Comment.Single: DocRGBColor(0x00, 0x80, 0x00),
        Token.Comment.Multiline: DocRGBColor(0x00, 0x80, 0x00),
        Token.Operator: DocRGBColor(0x3A, 0x3A, 0x3A),
        Token.Punctuation: DocRGBColor(0x3A, 0x3A, 0x3A),
        Token.Text: DocRGBColor(0x1A, 0x1A, 0x1A),
    }
    for key, color in color_map.items():
        if token_type in key:
            return color
    return DocRGBColor(0x1A, 0x1A, 0x1A)


async def call_deepseek(prompt_text: str) -> list[dict[str, Any]]:
    messages = [
        {"role": "system", "content": SLIDE_GENERATION_PROMPT},
        {"role": "user", "content": prompt_text},
    ]
    body = {
        "model": DEEPSEEK_MODEL,
        "messages": messages,
        "stream": False,
        "temperature": 0.7,
        "max_tokens": 8192,
    }
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
    }
    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(
            f"{DEEPSEEK_BASE_URL}/chat/completions",
            json=body,
            headers=headers,
        )
    if resp.status_code != 200:
        logger.error("DeepSeek API错误: %s %s", resp.status_code, resp.text[:500])
        raise HTTPException(status_code=502, detail=f"AI服务请求失败: {resp.status_code}")

    data = resp.json()
    content = data["choices"][0]["message"]["content"]
    logger.info("DeepSeek响应(前300字): %s", content[:300])

    try:
        obj = json.loads(extract_json(content))
        slides = obj.get("slides", [])
        if not slides:
            raise ValueError("AI未生成任何幻灯片")
        return slides
    except Exception as e:
        logger.error("JSON解析失败: %s", str(e))
        logger.error("原始内容(前800字): %s", content[:800])
        raise HTTPException(status_code=502, detail=f"AI返回格式异常: {str(e)}")


async def call_spark(prompt_text: str) -> list[dict[str, Any]]:
    messages = [
        {"role": "system", "content": SLIDE_GENERATION_PROMPT},
        {"role": "user", "content": prompt_text},
    ]
    body = {
        "model": SPARK_MODEL,
        "messages": messages,
        "stream": False,
        "temperature": 0.7,
        "max_tokens": 8192,
    }
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {SPARK_API_KEY}",
    }
    async with httpx.AsyncClient(timeout=120) as client:
        resp = await client.post(
            f"{SPARK_BASE_URL}/chat/completions",
            json=body,
            headers=headers,
        )
    if resp.status_code != 200:
        logger.error("Spark API错误: %s %s", resp.status_code, resp.text[:500])
        raise HTTPException(status_code=502, detail=f"Spark AI请求失败: {resp.status_code}")

    data = resp.json()
    content = data["choices"][0]["message"]["content"]
    logger.info("Spark响应(前300字): %s", content[:300])

    try:
        obj = json.loads(extract_json(content))
        slides = obj.get("slides", [])
        if not slides:
            raise ValueError("AI未生成任何幻灯片")
        return slides
    except Exception as e:
        logger.error("Spark JSON解析失败: %s", str(e))
        logger.error("原始内容(前800字): %s", content[:800])
        raise HTTPException(status_code=502, detail=f"Spark返回格式异常: {str(e)}")


def extract_json(text: str) -> str:
    text = text.strip()
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        json_str = text[start:end + 1]
        return _sanitize_json(json_str)
    return text


def _sanitize_json(json_str: str) -> str:
    sb = []
    i = 0
    while i < len(json_str):
        c = json_str[i]
        if c == '\\' and i + 1 < len(json_str):
            nxt = json_str[i + 1]
            if nxt in ('"', '\\', '/', 'b', 'f', 'n', 'r', 't', 'u'):
                sb.append(c)
            else:
                sb.append('\\\\')
        else:
            sb.append(c)
        i += 1
    return ''.join(sb)


def _inject_search_slide(slides_data: list[dict[str, Any]], search_text: str, item_count: str) -> list[dict[str, Any]]:
    raw_lines = [line.strip() for line in search_text.strip().split('\n') if line.strip()
                 if not line.strip().startswith('##')]
    items = []
    current_name = ""
    for line in raw_lines:
        if line.startswith('🔗'):
            url = line[2:].strip()
            if current_name and url:
                items.append({"name": current_name, "url": url})
                current_name = ""
        else:
            if current_name:
                items.append({"name": current_name, "url": ""})
            import re
            m = re.match(r'^\d+\.\s*(.+)', line)
            current_name = m.group(1).strip() if m else line.strip()
    if current_name:
        items.append({"name": current_name, "url": ""})

    count = int(item_count) if item_count else len([it for it in items if it["url"]])
    title = f"📚 全网搜索资源（共{count}条链接）" if count else "📚 全网搜索资源"

    return slides_data + [{
        "type": "search_results",
        "title": title,
        "items": items[:10]
    }]


def build_pptx(slides: list[dict[str, Any]], course_name: str, knowledge_point: str,
               images: list[bytes] | None = None) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_color = RGBColor(0x1A, 0x1A, 0x2E)
    body_color = RGBColor(0x33, 0x33, 0x33)
    accent_color = RGBColor(0x16, 0x5D, 0xFF)
    light_bg = RGBColor(0xF8, 0xF9, 0xFC)

    img_list = images or []
    img_index = 0

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            _add_title_slide(prs, slide_data, accent_color, light_bg, title_color)
        elif slide_type == "content":
            img_bytes = None
            if img_index < len(img_list):
                img_bytes = img_list[img_index]
                img_index += 1
            _add_content_slide(prs, slide_data, title_color, body_color, accent_color, light_bg, img_bytes)
        elif slide_type == "table":
            _add_table_slide(prs, slide_data, title_color, body_color, accent_color, light_bg)
        elif slide_type == "search_results":
            _add_search_results_slide(prs, slide_data, title_color, body_color, accent_color, light_bg)
        elif slide_type == "end":
            _add_end_slide(prs, slide_data, accent_color, light_bg, title_color)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs, data, accent_color, light_bg, title_color):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = title_color

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    author = data.get("author", "")

    if title:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.2))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.CENTER

    if author:
        txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.8))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = author
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x99, 0x99, 0xBB)
        p3.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs, data, title_color, body_color, accent_color, light_bg,
                       img_bytes: bytes | None = None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_title = data.get("title", "")
    bullets = data.get("bullets", [])

    header = slide.shapes.add_shape(
        1, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = light_bg
    header.line.fill.background()

    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.65), Inches(2.0), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    if bullets:
        if img_bytes:
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(6.5), Inches(4.8))
            img_stream = BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(8.0), Inches(2.2), Inches(4.8), Inches(4.5))
        else:
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.8))

        tf2 = txBox.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = Pt(22)
            p2.font.color.rgb = body_color
            p2.space_after = Pt(16)
            p2.level = 0
            p2.alignment = PP_ALIGN.LEFT


def _add_table_slide(prs, data, title_color, body_color, accent_color, light_bg):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide_title = data.get("title", "")
    headers = data.get("headers", [])
    rows = data.get("rows", [])

    header_shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = light_bg
    header_shape.line.fill.background()

    tf = header_shape.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = title_color

    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.65), Inches(2.0), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()

    num_rows = len(rows) + 1
    num_cols = len(headers) if headers else 1
    tbl_left = Inches(1.2)
    tbl_top = Inches(2.2)
    tbl_width = Inches(10.9)
    tbl_height = Inches(0.5 * num_rows)

    table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table

    for ci, head in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = head
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent_color

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            if ci < num_cols:
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.color.rgb = body_color
                    paragraph.alignment = PP_ALIGN.CENTER
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF6, 0xFA)


def _add_search_results_slide(prs, data, title_color, body_color, accent_color, light_bg):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)

    title = data.get("title", "搜索资源")
    items = data.get("items", [])

    header = slide.shapes.add_shape(1, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = light_bg
    header.line.fill.background()
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.LEFT

    icon_shape = slide.shapes.add_shape(1, Inches(10.5), Inches(0.5), Inches(1.8), Inches(1.0))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = accent_color
    icon_shape.line.fill.background()
    tf2 = icon_shape.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "🔍"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

    y_start = 2.0
    for i, item in enumerate(items[:8]):
        name = item.get("name", "") if isinstance(item, dict) else str(item)
        url = item.get("url", "") if isinstance(item, dict) else ""
        y = y_start + i * 0.7

        box = slide.shapes.add_shape(1, Inches(1.2), Inches(y), Inches(10.9), Inches(0.6))
        box.fill.solid()
        fill_color = light_bg if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        box.fill.fore_color.rgb = fill_color
        box.line.fill.background()

        tf3 = box.text_frame
        tf3.word_wrap = True
        tf3.margin_left = Inches(0.3)
        tf3.margin_top = Inches(0.05)

        p3 = tf3.paragraphs[0]
        p3.text = f"{i + 1}.  {name}"
        p3.font.size = Pt(16)
        p3.font.bold = True
        p3.font.color.rgb = body_color

        if url:
            p4 = tf3.add_paragraph()
            p4.text = f"     🔗 {url}"
            p4.font.size = Pt(12)
            p4.font.color.rgb = accent_color

    footer = slide.shapes.add_textbox(Inches(1.2), Inches(6.6), Inches(10.9), Inches(0.6))
    tf4 = footer.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "以上资源由讯飞ONE SEARCH智能搜索提供"
    p4.font.size = Pt(11)
    p4.font.italic = True
    p4.font.color.rgb = RGBColor(0x99, 0x99, 0xBB)
    p4.alignment = PP_ALIGN.LEFT


def _add_end_slide(prs, data, accent_color, light_bg, title_color):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = title_color

    end_title = data.get("title", "谢谢观看")
    end_subtitle = data.get("subtitle", "")

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = end_title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    if end_subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.0))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = end_subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=5050)
